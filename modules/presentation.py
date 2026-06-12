from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

class PresentationBuilder:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def add_slide(self, title: str, content_lines: list):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)

        left = Inches(0.8)
        top = Inches(0.5)
        width = Inches(11.7)
        height = Inches(1.2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.color.rgb = RGBColor(0x00, 0xd2, 0xff)
        p.font.bold = True

        left = Inches(0.8)
        top = Inches(1.8)
        width = Inches(11.7)
        height = Inches(5.0)
        txBox2 = slide.shapes.add_textbox(left, top, width, height)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        for i, line in enumerate(content_lines):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.text = line
            p2.font.size = Pt(18)
            p2.font.color.rgb = RGBColor(0xEC, 0xF0, 0xF1)
            p2.space_after = Pt(12)

        return slide

    def build(self, output_path: str):
        self.add_slide("Automated CV Screening Agent",
                       ["Solusi AI/ML untuk Otomatisasi Seleksi CV",
                        "",
                        "Allure Industries - Workshop Cibitung",
                        "AI Specialist Technical Assessment"])

        self.add_slide("Problem Statement",
                       ["• Proses seleksi CV manual oleh HR",
                        "  - Memakan waktu dan sumber daya besar",
                        "  - Rentan terhadap bias manusia",
                        "  - Tidak konsisten antar reviewer",
                        "",
                        "• Perusahaan dalam masa pertumbuhan dengan hiring intensitas tinggi",
                        "  - Volume CV besar setiap harinya",
                        "  - Kebutuhan respon cepat",
                        "",
                        "• Dampak:",
                        "  - Kandidat berkualitas terlewat",
                        "  - Time-to-hire panjang",
                        "  - Biaya rekrutmen tinggi"])

        self.add_slide("System Architecture Overview",
                       ["1. Input Layer",
                        "   • Job Portal Scraper (Jobstreet, Glints, dll)",
                        "   • Manual Upload (PDF/DOCX) via Web App",
                        "",
                        "2. Processing Layer",
                        "   • PDF Text Extraction (pdfplumber)",
                        "   • AI-Powered Information Extraction",
                        "   • Semantic Analysis (Sentence Transformers)",
                        "",
                        "3. Matching & Scoring Layer",
                        "   • Semantic Similarity (35%)",
                        "   • Skills Match (30%)",
                        "   • Keyword Relevance (20%)",
                        "   • Experience Score (15%)",
                        "",
                        "4. Output Layer",
                        "   • Ranked Candidate List",
                        "   • Google Sheets Integration",
                        "   • Separate CV Folder (Lolos)",
                        "   • Auto Email ke Tim HC"])

        self.add_slide("AI/ML Model Selection",
                       ["• Text Extraction: pdfplumber + Regex NLP",
                        "  - Ekstraksi nama, email, telepon, skill",
                        "",
                        "• Semantic Embedding: Sentence Transformers",
                        "  - Model: paraphrase-multilingual-MiniLM-L12-v2",
                        "  - Multilingual (support Bahasa Inggris & Indonesia)",
                        "  - Ringan, cepat, tanpa GPU",
                        "",
                        "• Matching Algorithm:",
                        "  - Cosine Similarity (semantic match)",
                        "  - Weighted Multi-Factor Scoring",
                        "  - Threshold-based filtering (>50%)",
                        "",
                        "• Future Enhancement:",
                        "  - Integrasi LLM (GPT/Claude/Gemini) untuk reasoning",
                        "  - LayoutLM untuk ekstraksi dari PDF complex"])

        self.add_slide("System Flow / Pipeline",
                       ["1. CV Collection",
                        "   → Scraper/RPA → Local Folder / Upload Web",
                        "",
                        "2. Preprocessing",
                        "   → PDF → Text → Clean → Structure",
                        "",
                        "3. Information Extraction",
                        "   → Name, Email, Skills, Education, Experience",
                        "",
                        "4. Matching & Scoring",
                        "   → Compare with Job Requirements",
                        "   → Weighted Semantic + Keyword + Skills + Exp",
                        "",
                        "5. Ranking & Filtering",
                        "   → Sort by Total Score → Filter >= 50%",
                        "",
                        "6. Output Generation",
                        "   → Google Sheets (candidate list)",
                        "   → Selected CVs folder",
                        "   → Email summary ke HC team"])

        self.add_slide("Automation & Scheduling",
                       ["• Daily Automation at 23:59",
                        "  - Otomatis proses semua CV baru",
                        "  - Generate hasil seleksi harian",
                        "  - Kirim email ringkasan ke tim HC",
                        "",
                        "• Web Scraper / RPA (Bonus)",
                        "  - Scrape Jobstreet & Glints secara otomatis",
                        "  - Download CV tanpa API resmi",
                        "  - Simpan ke folder input sistem",
                        "",
                        "• Integration Ready:",
                        "  - Google Sheets API",
                        "  - SMTP Email",
                        "  - Scheduled Tasks (Windows Task Scheduler)"])

        self.add_slide("Proof of Concept - Results",
                       ["• Testing dengan 5 CV Sample untuk posisi Junior Architect",
                        "",
                        "• Scoring Components:",
                        "  - Semantic Similarity: 35%",
                        "  - Skills Match: 30%",
                        "  - Keyword Relevance: 20%",
                        "  - Experience: 15%",
                        "",
                        "• Threshold Lolos: ≥ 50% Total Score",
                        "",
                        "• Output:",
                        "  - Ranked candidates with detailed scores",
                        "  - Alasan rekomendasi per kandidat",
                        "  - CV terpisah untuk yang lolos"])

        self.add_slide("Business Impact & Benefits",
                       ["• Efisiensi Waktu: 80% lebih cepat dari manual",
                        "  - 5 CV/detik vs 5-10 menit manual per CV",
                        "",
                        "• Objektivitas: Skor konsisten, zero bias",
                        "  - Kriteria seragam untuk semua kandidat",
                        "",
                        "• Skalabilitas: Bisa handle ratusan CV/hari",
                        "  - Tanpa perlu tambahan tenaga HR",
                        "",
                        "• Biaya: Minimal infrastructure (CPU only)",
                        "  - Open source stack, no license cost",
                        "",
                        "• Akurasi: Multi-factor scoring",
                        "  - Semantic + Keyword + Skills + Experience"])

        self.add_slide("Next Steps & Development Roadmap",
                       ["• Phase 1 (Hari 1-2): Core System ✅",
                        "  - CV parsing, extraction, matching, scoring",
                        "",
                        "• Phase 2 (Hari 3): Automation & Integration ✅",
                        "  - Google Sheets sync, email auto-send",
                        "  - Daily scheduler 23:59",
                        "",
                        "• Phase 3 (Minggu 2): Enhancement ✅",
                        "  - LLM integration (Gemini/OpenAI) untuk reasoning",
                        "  - Web scraper/RPA untuk Jobstreet & Glints",
                        "  - DOCX support, LinkedIn extraction",
                        "",
                        "• Phase 4 (Minggu 3-4): Production",
                        "  - Deployment ke server/cloud",
                        "  - Monitoring & dashboard"])

        self.add_slide("Analisis: Data Tambahan yang Diperlukan",
                       ["1. Database Referensi Skill & Industri",
                        "   - Standarisasi nama skill (misal: 'Ms. Office' vs 'Microsoft Office')",
                        "   - Mapping skill by industry/job family",
                        "",
                        "2. Historical Hiring Data",
                        "   - Data kandidat sebelumnya (lolos vs tidak)",
                        "   - Feedback dari interviewer untuk fine-tuning scoring",
                        "",
                        "3. Company Culture & Values Dictionary",
                        "   - Soft skills & cultural fit indicators",
                        "   - Keywords yang merepresentasikan budaya perusahaan",
                        "",
                        "4. Benchmark Gaji & Pasar",
                        "   - Rentang gaji per posisi & level",
                        "   - Ekspektasi kandidat vs budget perusahaan",
                        "",
                        "5. Database Sertifikasi & Pendidikan",
                        "   - Akreditasi institusi pendidikan",
                        "   - Validitas sertifikasi profesional"])

        self.add_slide("Analisis: Tools/Aplikasi Tambahan yang Dibutuhkan",
                       ["1. LLM API (Gemini/OpenAI/Claude)",
                        "   - Untuk extraction data CV yang lebih akurat",
                        "   - Reasoning & explanation generation",
                        "",
                        "2. Google Cloud / AWS Account",
                        "   - Google Sheets API + Service Account",
                        "   - Cloud deployment & auto-scaling",
                        "",
                        "3. Database (PostgreSQL / MongoDB)",
                        "   - Menyimpan history screening",
                        "   - Tracking kandidat dari waktu ke waktu",
                        "",
                        "4. Dashboard & Monitoring Tool",
                        "   - Grafana / Metabase untuk visualisasi",
                        "   - Track metrics: time-to-hire, acceptance rate",
                        "",
                        "5. Queue System (RabbitMQ / Redis)",
                        "   - Handle high-volume CV processing",
                        "   - Async task management",
                        "",
                        "6. Version Control & CI/CD",
                        "   - Git + GitHub/GitLab",
                        "   - Automated testing & deployment pipeline"])

        self.add_slide("System Architecture - Full Diagram",
                       ["INPUT LAYER:",
                        "  [Jobstreet] [Glints] [Manual Upload] [Email Inbox]",
                        "         |",
                        "         v",
                        "PROCESSING LAYER:",
                        "  [PDF Parser] [DOCX Parser] [AI Extractor] [OCR]",
                        "  [Text Cleaner] [Language Detector] [NER]",
                        "         |",
                        "         v",
                        "MATCHING & SCORING LAYER:",
                        "  [Semantic Embedding] [Skills Match] [Keyword Match]",
                        "  [Experience Calc] [Weighted Scoring] [LLM Reasoning]",
                        "         |",
                        "         v",
                        "OUTPUT LAYER:",
                        "  [Google Sheets] [Selected CVs Folder] [Email to HC]",
                        "  [Dashboard] [API Endpoints] [Notification]"])

        self.prs.save(output_path)
        return output_path
